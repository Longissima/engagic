"""Document extractor using PyMuPDF with OCR fallback

Moved from: infocore/processing/pdf_extractor.py

Supports: PDF, DOCX (via PyMuPDF), legacy .doc (via antiword), RTF (via striprtf),
PPTX (via python-pptx), XLSX (via openpyxl).

Legislative formatting detection (strikethrough/underline):
- Detects thin filled/stroked marks and PDF markup annotations
- Strikethrough = deletions from law (line through middle of text)
- Underline = additions to law (line below text)
- Outputs as [DELETED: text] and [ADDED: text] tags in markdown
"""

import io
import os
import re
import subprocess
import tempfile
import time
import warnings
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Dict, Any, List, Tuple, Optional, cast

import fitz  # PyMuPDF
import pytesseract
import requests
from PIL import Image

from config import get_logger
from exceptions import ExtractionError

logger = get_logger(__name__).bind(component="parser")

# Browser-like headers to avoid bot detection
DEFAULT_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36",
    "Accept": "application/pdf,application/octet-stream,*/*",
    "Accept-Language": "en-US,en;q=0.9",
    "Accept-Encoding": "gzip, deflate, br",
    "Connection": "keep-alive",
}

# Set conservative limit for OCR on scanned PDFs
# 100MP = ~300MB peak RAM (conservative for 2GB VPS with other services)
# Convert PIL warnings to errors to catch decompression bombs
Image.MAX_IMAGE_PIXELS = 100000000
warnings.simplefilter('error', Image.DecompressionBombWarning)

# Magic bytes for file format detection
_OLE2_MAGIC = b'\xd0\xcf\x11\xe0'    # Legacy .doc, .xls, .ppt (OLE2 Compound)
_ZIP_MAGIC = b'PK\x03\x04'           # .docx, .xlsx, .pptx (OOXML/ZIP)
_PDF_MAGIC = b'%PDF-'
_RTF_MAGIC = b'{\\rtf'

# PyMuPDF exposes these at runtime, but its type stubs omit them.
_PDF_ANNOT_UNDERLINE = cast(int, getattr(fitz, "PDF_ANNOT_UNDERLINE", 9))
_PDF_ANNOT_STRIKE_OUT = cast(int, getattr(fitz, "PDF_ANNOT_STRIKE_OUT", 11))

XLSX_MAX_ROWS_PER_SHEET = 100
XLSX_MAX_SHEETS = 20
XLSX_MAX_COLS = 30
XLSX_MAX_CELL_CHARS = 300
XLSX_MAX_TOTAL_CHARS = 200_000


def _detect_format(data: bytes) -> str:
    """Detect document format from magic bytes.

    Returns 'pdf', 'docx', 'pptx', 'xlsx', 'doc', 'rtf', or 'unknown'.
    For ZIP-based OOXML, peeks at [Content_Types].xml to distinguish
    PPTX (presentationml) from DOCX (wordprocessingml).
    """
    if data[:5] == _PDF_MAGIC:
        return "pdf"
    if data[:4] == _ZIP_MAGIC:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                ct = zf.read("[Content_Types].xml").decode("utf-8", errors="ignore")
                if "presentationml" in ct:
                    return "pptx"
                if "spreadsheetml" in ct:
                    return "xlsx"
        except Exception:
            pass
        return "docx"
    if data[:4] == _OLE2_MAGIC:
        return "doc"
    if data[:5] == _RTF_MAGIC:
        return "rtf"
    return "unknown"


def _extract_legacy_doc(data: bytes) -> Optional[str]:
    """Extract text from legacy .doc (OLE2) using antiword. Returns text or None."""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp_path = tmp.name
            tmp.write(data)
        result = subprocess.run(
            ["antiword", tmp_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
        logger.debug("antiword returned no text", exit_code=result.returncode, stderr=result.stderr[:200])
        return None
    except FileNotFoundError:
        logger.warning("antiword not installed, cannot extract legacy .doc")
        return None
    except subprocess.TimeoutExpired:
        logger.warning("antiword timed out")
        return None
    except Exception as e:
        logger.debug("antiword extraction failed", error=str(e))
        return None
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def _extract_rtf(data: bytes) -> Optional[str]:
    """Extract text from RTF using striprtf. Returns text or None."""
    try:
        from striprtf.striprtf import rtf_to_text
        rtf_str = data.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_str)
        return text if text and text.strip() else None
    except Exception as e:
        logger.debug("rtf extraction failed", error=str(e))
        return None


def _extract_pptx(data: bytes) -> Optional[str]:
    """Extract text from PPTX using python-pptx. Returns text or None."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                shape_obj = cast(Any, shape)
                if shape_obj.has_text_frame:
                    for para in shape_obj.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape_obj.has_table:
                    for row in shape_obj.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
        return "\n".join(parts) if parts else None
    except Exception as e:
        logger.debug("pptx extraction failed", error=str(e))
        return None


def _extract_xlsx(data: bytes) -> Optional[str]:
    """Extract XLSX as markdown tables using openpyxl. Returns text or None.

    Fee schedules and budget detail live in spreadsheets; without this, 800+
    attachments were invisible to summarization. Row cap keeps monster
    workbooks bounded -- the model needs the shape and leading rows, not every
    line of a 10,000-row ledger. Confidence: 8/10 on the row cap being enough.
    """
    wb = None
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        total_chars = 0
        size_limit_reached = False

        def append_part(value: str) -> bool:
            """Append within the workbook-wide output cap; return False if full."""
            nonlocal total_chars
            separator_chars = 1 if parts else 0
            remaining = XLSX_MAX_TOTAL_CHARS - total_chars - separator_chars
            if remaining <= 0:
                return False
            clipped = value[:remaining]
            parts.append(clipped)
            total_chars += separator_chars + len(clipped)
            return len(clipped) == len(value)

        for sheet_index, ws in enumerate(wb.worksheets):
            if sheet_index >= XLSX_MAX_SHEETS:
                append_part(
                    f"[{len(wb.worksheets) - XLSX_MAX_SHEETS} more sheets omitted]"
                )
                break
            if total_chars >= XLSX_MAX_TOTAL_CHARS:
                break
            rows_out = []
            total_rows = ws.max_row or 0
            rows_to_read = min(total_rows, XLSX_MAX_ROWS_PER_SHEET)
            for row in ws.iter_rows(
                min_row=1,
                max_row=rows_to_read,
                max_col=min(ws.max_column or 1, XLSX_MAX_COLS),
                values_only=True,
            ):
                cells = [
                    "" if v is None else str(v).strip()[:XLSX_MAX_CELL_CHARS]
                    for v in row[:XLSX_MAX_COLS]
                ]
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    rows_out.append("| " + " | ".join(cells) + " |")
            if rows_out:
                if not append_part(f"## Sheet: {ws.title}"):
                    size_limit_reached = True
                    break
                for line in rows_out:
                    if not append_part(line):
                        size_limit_reached = True
                        break
                if size_limit_reached:
                    break
                if total_rows > XLSX_MAX_ROWS_PER_SHEET:
                    if not append_part(
                        f"[{total_rows - XLSX_MAX_ROWS_PER_SHEET} more rows omitted]"
                    ):
                        size_limit_reached = True
                        break
        text = "\n".join(parts)
        if size_limit_reached:
            suffix = "\n\n[remaining workbook content omitted: extraction size limit]"
            text = text[: XLSX_MAX_TOTAL_CHARS - len(suffix)] + suffix
        return text if text.strip() else None
    except Exception as e:
        logger.debug("xlsx extraction failed", error=str(e))
        return None
    finally:
        if wb is not None:
            wb.close()


def _detect_horizontal_lines(page: fitz.Page) -> List[Tuple[float, float, float]]:
    """
    Detect horizontal lines from drawing instructions.

    Strikethrough/underline in MS Word/LibreOffice are rendered as THIN FILLED RECTANGLES.

    Returns:
        list of (x0, x1, y_position) tuples for horizontal bars
    """
    # Strikethrough is not one encoding. Known rendering families, each
    # detected independently and fed into the same mid-span classifier:
    #   A. Thin FILLED rects -- LibreOffice (black), Word track-changes (red).
    #   B. STROKED horizontal line/thin-rect paths -- other Word export paths
    #      and HTML-to-PDF vendors. (The old code nested 'l' handling inside
    #      the fill branch, so stroke-only paths could never match.)
    #   C. StrikeOut ANNOTATIONS -- PDF markup layer.
    # Near-white marks are excluded (masking rectangles). Geometry -- thin,
    # wide, crossing mid-span -- remains the real discriminator downstream.
    lines = []

    def _is_visible_color(color) -> bool:
        """Accept dark and saturated marks while rejecting white masks."""
        if color is None:
            return False
        components = color if isinstance(color, (tuple, list)) else (color,)
        return any(component < 0.9 for component in components[:3])

    def _append_thin_rect(rect) -> None:
        x0, y0, x1, y1 = rect
        height = abs(y1 - y0)
        width = abs(x1 - x0)
        if height < 2 and width > 5:
            lines.append((x0, x1, (y0 + y1) / 2))

    for path in page.get_drawings():
        fill_color = path.get("fill")
        stroke_color = path.get("color")
        visible_fill = _is_visible_color(fill_color)
        visible_stroke = _is_visible_color(stroke_color)
        if not (visible_fill or visible_stroke):
            continue
        for item in path["items"]:
            if item[0] == "re":
                _append_thin_rect(item[1])
            elif item[0] == "l" and visible_stroke:
                p1, p2 = item[1:3]
                if abs(p1.y - p2.y) < 0.5 and abs(p2.x - p1.x) > 5:
                    lines.append((min(p1.x, p2.x), max(p1.x, p2.x), (p1.y + p2.y) / 2))

    annot = page.first_annot
    while annot:
        annot_type = annot.type[0]
        if annot_type in (_PDF_ANNOT_STRIKE_OUT, _PDF_ANNOT_UNDERLINE):
            vertices = annot.vertices or []
            quads = [vertices[i : i + 4] for i in range(0, len(vertices), 4)]
            quads = [quad for quad in quads if len(quad) == 4]
            if quads:
                for quad in quads:
                    xs = [point[0] for point in quad]
                    ys = [point[1] for point in quad]
                    y = (
                        (min(ys) + max(ys)) / 2
                        if annot_type == _PDF_ANNOT_STRIKE_OUT
                        else max(ys)
                    )
                    if max(xs) - min(xs) > 5:
                        lines.append((min(xs), max(xs), y))
            else:
                rect = annot.rect
                y = (
                    (rect.y0 + rect.y1) / 2
                    if annot_type == _PDF_ANNOT_STRIKE_OUT
                    else rect.y1
                )
                lines.append((rect.x0, rect.x1, y))
        annot = annot.next

    return lines


def _span_mark_runs(
    chars: List[Tuple[str, float, float]],
    span_bbox: Tuple[float, float, float, float],
    lines: List[Tuple[float, float, float]],
) -> List[Tuple[str, Optional[str]]]:
    """Classify a span's characters against mark lines, character by character.

    A single span can carry BOTH a strike over its old words and an underline
    under its new words (Word emits "old new" as one span with two marks at
    different x-ranges). Span-granular classification either misses one mark
    or mislabels the whole span, so marks are applied per character using each
    line's x-extent, then grouped into runs.

    chars: (character, x0, x1) in reading order.
    Returns runs of (text, mark) where mark is None, 'del', or 'add'.
    """
    y0, y1 = span_bbox[1], span_bbox[3]
    height = y1 - y0
    strike_lo, strike_hi = y0 + 0.3 * height, y0 + 0.7 * height
    under_lo, under_hi = y1 - 1, y1 + 3

    marks: List[Optional[str]] = [None] * len(chars)
    for line_x0, line_x1, line_y in lines:
        if strike_lo <= line_y <= strike_hi:
            mark = "del"
        elif under_lo <= line_y <= under_hi:
            mark = "add"
        else:
            continue
        for i, (_, cx0, cx1) in enumerate(chars):
            mid = (cx0 + cx1) / 2
            if line_x0 - 1 <= mid <= line_x1 + 1:
                # Deletion outranks addition on the rare overlap; never
                # downgrade an already-marked deletion.
                if marks[i] is None or (mark == "del" and marks[i] == "add"):
                    marks[i] = mark

    runs: List[List] = []
    for (char, _, _), mark in zip(chars, marks):
        if runs and runs[-1][1] == mark:
            runs[-1][0] += char
        else:
            runs.append([char, mark])
    # Whitespace-only marked runs are geometric accidents (a mark line
    # skimming a space between words); render them plain.
    return [
        (text, None if (mark and not text.strip()) else mark)
        for text, mark in runs
    ]


def _page_mark_runs(
    page: fitz.Page,
    detected_lines: Optional[List[Tuple[float, float, float]]] = None,
) -> Tuple[List[List[List[Tuple[str, Optional[str]]]]], int, int]:
    """Compute per-span mark runs for a whole page.

    Returns (blocks -> lines -> marked runs in reading order,
    struck_span_count, underlined_span_count) so extraction and the activation
    gate share one deterministic classification pass without flattening the
    document's line and paragraph structure.
    """
    lines = detected_lines if detected_lines is not None else _detect_horizontal_lines(page)
    page_runs: List[List[List[Tuple[str, Optional[str]]]]] = []
    struck = 0
    underlined = 0
    raw = cast(
        Dict[str, Any],
        page.get_text("rawdict", sort=True),  # type: ignore[attr-defined]
    )
    for block in raw["blocks"]:
        if block.get("type") != 0:
            continue
        block_runs: List[List[Tuple[str, Optional[str]]]] = []
        for line_obj in block.get("lines", []):
            line_runs: List[Tuple[str, Optional[str]]] = []
            for span in line_obj.get("spans", []):
                chars = [
                    (ch["c"], ch["bbox"][0], ch["bbox"][2])
                    for ch in span.get("chars", [])
                ]
                if not chars:
                    continue
                runs = (
                    _span_mark_runs(chars, tuple(span["bbox"]), lines)
                    if lines
                    else [("".join(c for c, _, _ in chars), None)]
                )
                line_runs.extend(runs)
                # A mark must cover at least 2 non-space characters to count
                # as evidence -- single-character grazes are noise.
                if any(m == "del" and len(t.strip()) >= 2 for t, m in runs):
                    struck += 1
                if any(m == "add" and len(t.strip()) >= 2 for t, m in runs):
                    underlined += 1
            if line_runs:
                block_runs.append(line_runs)
        if block_runs:
            page_runs.append(block_runs)
    return page_runs, struck, underlined


def _has_legislative_legend(doc: fitz.Document, proximity_chars: int = 200, max_pages: int = 5) -> bool:
    """Check if document contains legislative formatting legend (clustered keywords).

    A true legislative legend has all 4 keyword types appearing close together,
    like: "Additions shown as underline, deletions shown as strikethrough"

    False positives occur when keywords are scattered throughout ordinance text
    describing what amendments do (e.g., "the addition of Section 12-345...").

    Legislative legends typically appear on the first few pages (cover or TOC),
    so we limit search to max_pages for performance.

    Args:
        doc: PyMuPDF document
        proximity_chars: Maximum distance between keywords to consider clustered (default 200)
        max_pages: Maximum pages to search (default 5 - legends are in first few pages)

    Returns:
        True only if all 4 keyword types appear within proximity_chars of each other
    """
    # Keyword patterns for each category
    addition_pattern = re.compile(r'\b(addition|added)\b')
    deletion_pattern = re.compile(r'\b(deletion|deleted)\b')
    underline_pattern = re.compile(r'\bunderline\b')
    strikethrough_pattern = re.compile(r'\bstrikethrough\b')

    # Search only first max_pages (legends appear early in documents)
    pages_to_search = min(len(doc), max_pages)
    for page_num in range(pages_to_search):
        text = doc[page_num].get_text().lower()  # type: ignore[attr-defined]

        # Find all positions of each keyword type
        addition_positions = [m.start() for m in addition_pattern.finditer(text)]
        deletion_positions = [m.start() for m in deletion_pattern.finditer(text)]
        underline_positions = [m.start() for m in underline_pattern.finditer(text)]
        strikethrough_positions = [m.start() for m in strikethrough_pattern.finditer(text)]

        # All 4 types must be present on this page
        if not (addition_positions and deletion_positions and
                underline_positions and strikethrough_positions):
            continue

        # Check if any combination of positions clusters within proximity_chars
        # Use underline as anchor since it's least likely to appear in normal text
        for u_pos in underline_positions:
            window_start = u_pos - proximity_chars
            window_end = u_pos + proximity_chars

            has_addition = any(window_start <= p <= window_end for p in addition_positions)
            has_deletion = any(window_start <= p <= window_end for p in deletion_positions)
            has_strikethrough = any(window_start <= p <= window_end for p in strikethrough_positions)

            if has_addition and has_deletion and has_strikethrough:
                return True

    return False


# Three independent strikes are strong evidence on their own. A paired
# deletion/addition is also sufficient: compact amendments often change only
# one phrase, while ordinary tables and hyperlinks rarely produce both marks.
REDLINE_ACTIVATION_STRUCK_SPANS = 3


def _redline_evidence_activates(struck_spans: int, underlined_spans: int) -> bool:
    return struck_spans >= REDLINE_ACTIVATION_STRUCK_SPANS or (
        struck_spans >= 1 and underlined_spans >= 1
    )


def count_redline_evidence(doc: fitz.Document, max_pages: int = 30) -> Dict[str, int]:
    """Count geometric strikethrough/underline evidence across the document.

    Deterministic activation signal for legislative formatting: a thin filled
    rectangle crossing the vertical MIDDLE of a text span (30-70% of span
    height, per _span_mark_runs) is a strikethrough candidate. Table
    rules and separators sit between spans and underlines sit at the
    baseline, so mid-span crossings are the discriminator that survives the
    false positives which motivated the old legend-only gate.

    Scans up to max_pages pages, stopping early once enough strike evidence
    accumulates to activate. Returns counts for logging/auditing:
        {'struck_spans': N, 'underlined_spans': M, 'pages_scanned': P}
    """
    struck = 0
    underlined = 0
    pages_scanned = 0
    for page_num in range(min(len(doc), max_pages)):
        pages_scanned += 1
        lines = _detect_horizontal_lines(doc[page_num])
        if not lines:
            continue
        _, page_struck, page_underlined = _page_mark_runs(
            doc[page_num], detected_lines=lines
        )
        struck += page_struck
        underlined += page_underlined
        if _redline_evidence_activates(struck, underlined):
            break
    return {
        "struck_spans": struck,
        "underlined_spans": underlined,
        "pages_scanned": pages_scanned,
    }


def _extract_text_with_formatting(page: fitz.Page, page_num: int) -> str:
    """
    Extract text from page with legislative formatting tags.

    Returns text with [DELETED: ...] and [ADDED: ...] tags rendered at
    character-run granularity: a single span holding both struck old text and
    underlined new text yields separate tagged runs, not one mislabeled span.
    """
    lines = _detect_horizontal_lines(page)
    if not lines:
        # No mark lines on this page, return plain text
        return cast(str, page.get_text(sort=True))  # type: ignore[attr-defined]

    page_runs, _, _ = _page_mark_runs(page, detected_lines=lines)

    def render_line(runs: List[Tuple[str, Optional[str]]]) -> str:
        parts = []
        for text, mark in runs:
            if mark == "del":
                parts.append(f"[DELETED: {text}]")
            elif mark == "add":
                parts.append(f"[ADDED: {text}]")
            else:
                parts.append(text)
        return "".join(parts)

    return "\n\n".join(
        "\n".join(render_line(line_runs) for line_runs in block_runs)
        for block_runs in page_runs
    )


class PdfExtractor:
    """PDF extractor using PyMuPDF with OCR fallback"""

    _instance: "PdfExtractor | None" = None

    @classmethod
    def shared(cls) -> "PdfExtractor":
        """Return a shared singleton instance."""
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    def __init__(self, ocr_threshold: int = 100, ocr_dpi: int = 200, detect_legislative_formatting: bool = True, max_ocr_workers: int | None = None, ocr_enabled: bool = True):
        """Initialize PDF extractor

        Args:
            ocr_threshold: Minimum characters per page before triggering OCR fallback
            ocr_dpi: DPI for image rendering when using OCR (higher = better quality, slower)
                    Default 200 - balances quality vs memory on multi-core VPS
            detect_legislative_formatting: If True, detect strikethrough (deletions) and underline (additions)
                    in legislative documents and tag them as [DELETED: ...] and [ADDED: ...].
                    Activates from a formatting legend or sufficient geometric mark evidence.
                    Default: True.
            max_ocr_workers: Maximum parallel OCR workers. Default: CPU count (min 1, max 4).
                    OCR is CPU-bound so more workers than cores causes thrashing.
            ocr_enabled: If False, pages below ocr_threshold keep their (thin) text-layer
                    text and are COUNTED in the result's ocr_pending instead of OCR'd.
                    This is the sync-side ground-truth mode (docs/CORPUS_ARCHITECTURE.md):
                    ocr_pending == 0 proves the output is identical to what the OCR-enabled
                    extractor would produce, so it is safe to persist to the corpus;
                    ocr_pending > 0 means the document belongs to the OCR-owning path.
        """
        self.ocr_threshold = ocr_threshold
        self.ocr_dpi = ocr_dpi
        self.ocr_enabled = ocr_enabled
        self.detect_legislative_formatting = detect_legislative_formatting

        # Auto-detect optimal worker count based on CPU cores
        # OCR is CPU-bound, so workers > cores = context switch overhead
        if max_ocr_workers is None:
            cpu_count = os.cpu_count() or 1
            self.max_ocr_workers = min(cpu_count, 4)  # Cap at 4 for memory
        else:
            self.max_ocr_workers = max_ocr_workers

        # Prevent Tesseract internal threading when running multiple workers
        # Each worker gets 1 thread to avoid CPU thrashing
        os.environ.setdefault('OMP_THREAD_LIMIT', '1')

        logger.info(
            "PDF extractor initialized",
            ocr_workers=self.max_ocr_workers,
            ocr_dpi=self.ocr_dpi,
            cpu_count=os.cpu_count()
        )

    def _render_page_for_ocr(self, page) -> Optional[Tuple[bytes, int, int]]:
        """Render page to PNG bytes for OCR (main thread, not thread-safe)

        Args:
            page: PyMuPDF page object

        Returns:
            Tuple of (png_bytes, width, height) or None if too large
        """
        try:
            pix = page.get_pixmap(dpi=self.ocr_dpi)
            megapixels = (pix.width * pix.height) / 1000000

            logger.debug(
                "rendering page for OCR",
                page_num=page.number + 1,
                dpi=self.ocr_dpi,
                width=pix.width,
                height=pix.height,
                megapixels=round(megapixels, 1)
            )

            if megapixels > 100:
                logger.warning(
                    "page image too large for OCR",
                    page_num=page.number + 1,
                    megapixels=round(megapixels, 1)
                )
                return None

            return (pix.tobytes("png"), pix.width, pix.height)
        except Exception as e:  # Intentionally broad: graceful degradation, returns None
            logger.error("failed to render page", page_num=page.number + 1, error=str(e))
            return None

    def _ocr_from_bytes(self, png_bytes: bytes, page_num: int) -> str:
        """Run OCR on pre-rendered PNG bytes (thread-safe)

        Args:
            png_bytes: PNG image bytes
            page_num: Page number (1-indexed, for logging)

        Returns:
            Extracted text from OCR, or empty string on failure
        """
        try:
            img = Image.open(io.BytesIO(png_bytes)).convert('L')  # Grayscale
            # --oem 1: LSTM-only (2-3x faster than legacy+LSTM default)
            # --psm 3: Fully automatic page segmentation
            # timeout: Hard 60s cap per page to prevent hangs
            text = pytesseract.image_to_string(
                img,
                config='--oem 1 --psm 3',
                timeout=60
            )
            return text
        except (Image.DecompressionBombError, Image.DecompressionBombWarning):
            logger.warning("page image too large for OCR", page_num=page_num)
            return ""
        except RuntimeError as e:
            # pytesseract raises RuntimeError on timeout
            if "Tesseract process timeout" in str(e):
                logger.warning("OCR timeout on page", page_num=page_num)
            else:
                logger.error("OCR runtime error", page_num=page_num, error=str(e))
            return ""
        except (OSError, pytesseract.TesseractError) as e:
            logger.error("OCR failed", page_num=page_num, error=str(e), error_type=type(e).__name__)
            return ""

    def _ocr_page(self, page) -> str:
        """Extract text from page using OCR (sequential fallback)

        Args:
            page: PyMuPDF page object

        Returns:
            Extracted text from OCR, or empty string if image too large
        """
        rendered = self._render_page_for_ocr(page)
        if rendered is None:
            return ""
        png_bytes, _, _ = rendered
        return self._ocr_from_bytes(png_bytes, page.number + 1)

    def _ocr_pages_parallel(self, ocr_tasks: List[Tuple[int, bytes, str]]) -> Dict[int, str]:
        """Run OCR on multiple pages in parallel

        Args:
            ocr_tasks: List of (page_num, png_bytes, original_text) tuples

        Returns:
            Dict mapping page_num to OCR result text
        """
        if not ocr_tasks:
            return {}

        results = {}
        workers = min(self.max_ocr_workers, len(ocr_tasks))

        logger.info(
            "starting parallel OCR",
            pages=len(ocr_tasks),
            workers=workers
        )

        ocr_start = time.time()

        # Total OCR budget: 5 minutes for all pages combined
        # Prevents infinite hangs if Tesseract locks up on any page
        ocr_total_timeout = 300

        with ThreadPoolExecutor(max_workers=workers) as executor:
            # Submit all OCR jobs
            future_to_page = {
                executor.submit(self._ocr_from_bytes, png_bytes, page_num): (page_num, original_text)
                for page_num, png_bytes, original_text in ocr_tasks
            }

            # Collect results as they complete (with total timeout)
            completed_futures = set()
            try:
                for future in as_completed(future_to_page, timeout=ocr_total_timeout):
                    completed_futures.add(future)
                    page_num, original_text = future_to_page[future]
                    try:
                        # Per-page timeout as secondary safeguard
                        ocr_result = future.result(timeout=120)

                        # Decide whether to use OCR or keep original
                        if self._is_ocr_better(original_text, ocr_result, page_num):
                            results[page_num] = ocr_result
                        else:
                            results[page_num] = original_text

                    except TimeoutError:
                        logger.warning("OCR timeout on page", page_num=page_num)
                        results[page_num] = original_text
                    except Exception as e:  # Intentionally broad: catch any thread exception
                        logger.error("parallel OCR failed", page_num=page_num, error=str(e) or type(e).__name__)
                        results[page_num] = original_text

            except TimeoutError:
                # Total OCR budget exceeded - cancel remaining futures and use original text
                timed_out_pages = []
                for future, (page_num, original_text) in future_to_page.items():
                    if future not in completed_futures:
                        future.cancel()
                        results[page_num] = original_text
                        timed_out_pages.append(page_num)
                logger.warning(
                    "OCR total timeout exceeded, skipping remaining pages",
                    timed_out_pages=timed_out_pages,
                    timeout_seconds=ocr_total_timeout
                )

        ocr_time = time.time() - ocr_start
        logger.info(
            "parallel OCR complete",
            pages=len(ocr_tasks),
            ocr_time=round(ocr_time, 2),
            avg_per_page=round(ocr_time / len(ocr_tasks), 2) if ocr_tasks else 0
        )

        return results

    def _is_ocr_better(self, original: str, ocr_result: str, page_num: int) -> bool:
        """Determine if OCR result is better than original text

        Args:
            original: Original extracted text
            ocr_result: OCR-produced text
            page_num: Page number (1-indexed, for logging)

        Returns:
            True if OCR should be used, False if original should be kept
        """
        orig_chars = len(original.strip())
        ocr_chars = len(ocr_result.strip())

        # If OCR produced nothing, always keep original
        if ocr_chars == 0:
            logger.info("keeping original text, OCR produced nothing", page_num=page_num)
            return False

        # Calculate quality metrics for OCR
        ocr_letters = sum(1 for c in ocr_result if c.isalpha())
        ocr_letter_ratio = ocr_letters / len(ocr_result) if len(ocr_result) > 0 else 0
        ocr_words = len(ocr_result.split())

        # OCR is better if:
        # 1. Original was empty -- keep any OCR output that isn't pure noise.
        #    Warrant registers, check runs, and tax rolls are numeric-heavy
        #    and dip below the 40% letter-ratio floor; dropping them would
        #    lose the page entirely.
        # 2. OR produced significantly more text (2x+ chars) at >40% letters
        # 3. OR produced more text at >70% letters
        empty_original_ocr_usable = orig_chars == 0 and ocr_chars >= 100 and ocr_letter_ratio > 0.15
        significantly_more = ocr_chars >= (orig_chars * 2) and ocr_letter_ratio > 0.4
        high_quality_improvement = ocr_chars > orig_chars and ocr_letter_ratio > 0.7

        if empty_original_ocr_usable or significantly_more or high_quality_improvement:
            logger.info(
                f"[PyMuPDF] Page {page_num}: Using OCR "
                f"({ocr_chars} chars, {ocr_words} words, {ocr_letter_ratio:.1%} letters > original {orig_chars} chars)"
            )
            return True
        else:
            logger.info(
                f"[PyMuPDF] Page {page_num}: Keeping original "
                f"(OCR: {ocr_chars} chars, {ocr_letter_ratio:.1%} letters not better than original {orig_chars} chars)"
            )
            return False

    def _extract_from_document(self, doc: fitz.Document, extract_links: bool, start_time: float) -> Dict[str, Any]:
        """Core extraction logic for opened PDF document

        Args:
            doc: Opened PyMuPDF document
            extract_links: Whether to extract hyperlinks
            start_time: Extraction start time (for timing)

        Returns:
            Dict with extraction results
        """
        page_texts = {}  # page_num -> text
        all_links = []
        ocr_tasks = []  # List of (page_num, png_bytes, original_text)
        ocr_pending = 0  # below-threshold pages skipped because ocr_enabled=False

        # Activation is deterministic and auditable: a textual legend OR
        # geometric strike evidence (spans crossed mid-height by strike
        # rects). The legend-only gate missed every real-world redline that
        # lacked a legend sentence -- deletions then extracted as plain text,
        # indistinguishable from operative language.
        use_formatting = False
        if self.detect_legislative_formatting:
            legend = _has_legislative_legend(doc)
            evidence = {"struck_spans": 0, "underlined_spans": 0, "pages_scanned": 0}
            if not legend:
                evidence = count_redline_evidence(doc)
            use_formatting = legend or _redline_evidence_activates(
                evidence["struck_spans"], evidence["underlined_spans"]
            )
            if use_formatting:
                logger.info(
                    "[PyMuPDF] Legislative formatting detected - tagging additions/deletions",
                    activated_by="legend" if legend else "geometry",
                    struck_spans=evidence["struck_spans"],
                    underlined_spans=evidence["underlined_spans"],
                )
            elif evidence["struck_spans"] > 0:
                logger.info(
                    "[PyMuPDF] Strike marks below activation threshold, not tagging",
                    struck_spans=evidence["struck_spans"],
                    threshold=REDLINE_ACTIVATION_STRUCK_SPANS,
                )

        # Per-page sanity cap. A single PDF page should never yield more than
        # a few KB of real text. When PyMuPDF sees a broken font CMap or a
        # mis-tagged content stream, get_text() can return megabytes of
        # gibberish for one page and the total exceeds the 1GB RLIMIT_AS
        # child budget or blows up multiprocessing pickle. Truncate early so
        # the rest of the document still extracts.
        _MAX_PAGE_CHARS = 200_000

        # Pass 1: Extract text from all pages, collect OCR tasks
        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text (with or without legislative formatting detection)
            if use_formatting:
                page_text = _extract_text_with_formatting(page, page_num + 1)
            else:
                page_text = cast(
                    str,
                    page.get_text(sort=True),  # type: ignore[attr-defined]
                )

            if len(page_text) > _MAX_PAGE_CHARS:
                logger.warning(
                    "[PyMuPDF] Page yielded suspicious text volume, truncating",
                    page_num=page_num + 1,
                    chars=len(page_text),
                    limit=_MAX_PAGE_CHARS,
                )
                page_text = page_text[:_MAX_PAGE_CHARS]

            initial_char_count = len(page_text.strip())

            # If page has minimal text, queue for OCR
            if initial_char_count < self.ocr_threshold:
                if not self.ocr_enabled:
                    # Ground-truth mode: count instead of OCR. The caller uses
                    # ocr_pending to decide whether this text is complete.
                    ocr_pending += 1
                    page_texts[page_num + 1] = page_text
                else:
                    logger.debug(
                        "page queued for OCR",
                        page_num=page_num + 1,
                        char_count=initial_char_count,
                        threshold=self.ocr_threshold
                    )
                    # Render page to PNG (main thread, PyMuPDF not thread-safe)
                    rendered = self._render_page_for_ocr(page)
                    if rendered:
                        png_bytes, _, _ = rendered
                        ocr_tasks.append((page_num + 1, png_bytes, page_text))
                    else:
                        # Rendering failed, keep original
                        page_texts[page_num + 1] = page_text
            else:
                page_texts[page_num + 1] = page_text

            # Extract links if requested
            if extract_links:
                page_links = page.get_links()  # type: ignore[attr-defined]
                for link in page_links:
                    if 'uri' in link and link['uri']:
                        all_links.append({
                            'page': page_num + 1,
                            'url': link['uri'],
                            'rect': link.get('from', None),
                        })

        page_count = len(doc)

        # Pass 2: Run OCR in parallel (outside doc context, PNG bytes already captured)
        if ocr_tasks:
            ocr_results = self._ocr_pages_parallel(ocr_tasks)
            page_texts.update(ocr_results)

        # Count OCR pages (pages where OCR was actually used, not just attempted)
        ocr_pages = sum(
            1 for page_num, _, original in ocr_tasks
            if page_num in page_texts and page_texts[page_num] != original
        )

        # Assemble final text in page order
        text_parts = [
            f"--- PAGE {page_num} ---\n{page_texts[page_num]}"
            for page_num in sorted(page_texts.keys())
        ]
        full_text = "\n\n".join(text_parts)

        extraction_time = time.time() - start_time

        # Determine extraction method
        method = "pymupdf+ocr" if ocr_pages > 0 else "pymupdf"

        log_msg = f"[PyMuPDF] Extracted {page_count} pages, {len(full_text)} chars"
        if ocr_pages > 0:
            log_msg += f" (OCR: {ocr_pages} pages)"
        if extract_links:
            log_msg += f", {len(all_links)} links"
        log_msg += f" in {extraction_time:.2f}s"
        logger.info(log_msg)

        result = {
            "success": True,
            "text": full_text,
            "method": method,
            "page_count": page_count,
            "extraction_time": extraction_time,
            "ocr_pages": ocr_pages,
            "ocr_pending": ocr_pending,
        }

        if extract_links:
            result["links"] = all_links

        return result

    def extract_from_path(self, pdf_path: str, extract_links: bool = False) -> Dict[str, Any]:
        """Extract text from a PDF on disk without loading bytes into memory.

        PDF-only (no format sniffing): fitz maps the file directly, so this
        is the memory-lean entry for guarded subprocess children that already
        hold a temp file. Mixed formats (docx/rtf/pptx) go through
        extract_from_bytes.
        """
        start_time = time.time()
        try:
            with fitz.open(pdf_path) as doc:
                return self._extract_from_document(doc, extract_links, start_time)
        except Exception as e:  # Intentionally broad: API boundary, convert to typed error
            extraction_time = time.time() - start_time
            logger.error("[PyMuPDF] extraction from path failed", error=str(e), error_type=type(e).__name__, extraction_time=round(extraction_time, 2))
            raise ExtractionError(
                f"PDF extraction failed after {extraction_time:.1f}s",
                document_type="pdf",
                original_error=e
            ) from e

    def extract_from_url(self, url: str, extract_links: bool = False) -> Dict[str, Any]:
        """Extract text and optionally links from PDF URL

        Args:
            url: PDF URL to extract from
            extract_links: Whether to extract hyperlinks (default False for backward compatibility)

        Returns dict with extraction results:
        {
            'success': bool,
            'text': str,
            'method': str,
            'page_count': int,
            'extraction_time': float,
            'links': list (if extract_links=True),
            'error': str (if failed)
        }
        """
        start_time = time.time()

        try:
            response = requests.get(url, headers=DEFAULT_HEADERS, timeout=30)
            response.raise_for_status()
            pdf_bytes = response.content

            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                return self._extract_from_document(doc, extract_links, start_time)

        except Exception as e:  # Intentionally broad: API boundary, convert to typed error
            extraction_time = time.time() - start_time
            logger.error("[PyMuPDF] extraction failed", url=url[:100], error=str(e), error_type=type(e).__name__, extraction_time=round(extraction_time, 2))
            raise ExtractionError(
                f"PDF extraction failed after {extraction_time:.1f}s",
                document_url=url,
                document_type="pdf",
                original_error=e
            ) from e

    def extract_from_bytes(self, pdf_bytes: bytes, extract_links: bool = False) -> Dict[str, Any]:
        """Extract text from document bytes (PDF, Office documents, or RTF).

        Detects format from magic bytes and routes to the appropriate extractor.
        PDF and DOCX go through PyMuPDF; legacy .doc uses antiword; RTF uses
        striprtf; PPTX and XLSX use their corresponding Office libraries.

        Args:
            pdf_bytes: Document content as bytes
            extract_links: If True, also extract hyperlinks (PDF/DOCX only)

        Returns dict with extraction results (same format as extract_from_url)
        """
        start_time = time.time()
        fmt = _detect_format(pdf_bytes)

        # Legacy .doc (OLE2) -- fitz can't handle this format
        if fmt == "doc":
            text = _extract_legacy_doc(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted legacy .doc via antiword", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "antiword",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("Legacy .doc extraction failed (antiword unavailable or returned no text)", document_type="doc")

        # RTF
        if fmt == "rtf":
            text = _extract_rtf(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted rtf", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "striprtf",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("RTF extraction failed", document_type="rtf")

        # PPTX
        if fmt == "pptx":
            text = _extract_pptx(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted pptx via python-pptx", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "python-pptx",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("PPTX extraction failed", document_type="pptx")

        # XLSX
        if fmt == "xlsx":
            text = _extract_xlsx(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted xlsx via openpyxl", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "openpyxl",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("XLSX extraction failed", document_type="xlsx")

        # PDF, DOCX, and unknown formats -- PyMuPDF handles all of these
        try:
            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                return self._extract_from_document(doc, extract_links, start_time)

        except Exception as e:  # Intentionally broad: API boundary, convert to typed error
            extraction_time = time.time() - start_time
            logger.error("[PyMuPDF] extraction from bytes failed", format=fmt, error=str(e), error_type=type(e).__name__, extraction_time=round(extraction_time, 2))
            raise ExtractionError(
                f"Document extraction failed after {extraction_time:.1f}s",
                document_type=fmt,
                original_error=e
            ) from e

    def validate_text(self, text: str) -> bool:
        """Validate text quality - basic check for now"""
        # Simple validation: check if text is not empty and has reasonable content
        if not text or len(text) < 100:
            return False

        # Check if text has reasonable letter ratio
        letters = sum(1 for c in text if c.isalpha())
        if letters / len(text) < 0.3:
            return False

        return True


def extract_document_file(pdf_path: str, ocr_threshold: int, ocr_dpi: int,
                          detect_legislative_formatting: bool, max_ocr_workers: int | None) -> Dict[str, Any]:
    """Guard-child target for process extraction (parsing.subprocess_guard).

    Module-level on purpose: forkserver children import the target's module,
    and this one costs a parsing-stack import instead of the analyzer's
    HTTP/LLM stack. Reads the document from disk (parent wrote a temp file
    and released the bytes) and routes through extract_from_bytes for
    format sniffing -- process attachments arrive as PDF/DOCX/RTF/PPTX/XLSX.
    """
    with open(pdf_path, "rb") as f:
        data = f.read()
    extractor = PdfExtractor(ocr_threshold, ocr_dpi, detect_legislative_formatting, max_ocr_workers)
    return extractor.extract_from_bytes(data)
